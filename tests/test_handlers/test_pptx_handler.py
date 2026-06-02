"""Tests for PptxHandler — PowerPoint text extraction."""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from core.result import Failure, Success
from handlers.pptx_handler import PptxHandler


@pytest.fixture
def pptx_with_title_and_body(tmp_path: Path) -> Path:
    """Single slide with a title (shape_type=13) and a text box body."""
    path = tmp_path / "titled.pptx"
    prs = Presentation()
    # Blank layout is index 6 in the default template
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add a title placeholder manually as a text box — shape_type 13 = TITLE
    # We use add_textbox and forcibly mark it via the spec — but python-pptx
    # doesn't expose shape_type setter. Instead use a title layout.
    title_layout = prs.slide_layouts[0]  # "Title Slide" layout
    slide2 = prs.slides.add_slide(title_layout)
    title_shape = slide2.shapes.title
    if title_shape is not None:
        title_shape.text = "My Slide Title"
    # Add a body text box
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    txBox.text_frame.text = "Body content here"
    prs.save(str(path))
    return path


@pytest.fixture
def pptx_no_title(tmp_path: Path) -> Path:
    """Single slide with only a text box (no title shape)."""
    path = tmp_path / "notitle.pptx"
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = "No title here"
    prs.save(str(path))
    return path


@pytest.fixture
def empty_pptx(tmp_path: Path) -> Path:
    """Presentation with no slides."""
    path = tmp_path / "empty.pptx"
    prs = Presentation()
    prs.save(str(path))
    return path


class TestPptxHandlerCanHandle:
    def test_lowercase_pptx(self) -> None:
        assert PptxHandler().can_handle(Path("file.pptx")) is True

    def test_uppercase_pptx(self) -> None:
        assert PptxHandler().can_handle(Path("file.PPTX")) is True

    def test_ppt_not_handled(self) -> None:
        assert PptxHandler().can_handle(Path("file.ppt")) is False

    def test_docx_not_handled(self) -> None:
        assert PptxHandler().can_handle(Path("file.docx")) is False


class TestPptxHandlerExtract:
    def test_slide_without_title_has_plain_header(self, pptx_no_title: Path) -> None:
        result = PptxHandler().extract(pptx_no_title)
        assert isinstance(result, Success)
        assert "[Slide 1]" in result.value.text

    def test_slide_body_text_present(self, pptx_no_title: Path) -> None:
        result = PptxHandler().extract(pptx_no_title)
        assert isinstance(result, Success)
        assert "No title here" in result.value.text

    def test_is_md_false(self, pptx_no_title: Path) -> None:
        result = PptxHandler().extract(pptx_no_title)
        assert isinstance(result, Success)
        assert result.value.is_md is False

    def test_source_path(self, pptx_no_title: Path) -> None:
        result = PptxHandler().extract(pptx_no_title)
        assert isinstance(result, Success)
        assert result.value.source_path == pptx_no_title

    def test_empty_presentation_returns_success_empty_text(
        self, empty_pptx: Path
    ) -> None:
        result = PptxHandler().extract(empty_pptx)
        assert isinstance(result, Success)
        assert result.value.text == ""

    def test_nonexistent_path_returns_failure(self, tmp_path: Path) -> None:
        result = PptxHandler().extract(tmp_path / "ghost.pptx")
        assert isinstance(result, Failure)
        assert result.recoverable is False

    def test_slide_with_title_placeholder_extracted(
        self, pptx_with_title_and_body: Path
    ) -> None:
        result = PptxHandler().extract(pptx_with_title_and_body)
        assert isinstance(result, Success)
        assert '[Slide' in result.value.text
        assert '"My Slide Title"' in result.value.text
