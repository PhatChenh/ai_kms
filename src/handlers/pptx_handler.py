"""PPTX handler — extracts slide text from PowerPoint files via python-pptx.

Files larger than CONFIG.main.handlers.max_file_size_bytes are rejected before
any parse work to prevent a single large drop from OOMing the capture process.
"""
from __future__ import annotations

from pathlib import Path

import structlog

from core.result import Failure, Result, Success
from handlers.base import BaseHandler, RawContent
from handlers.registry import HandlerRegistry

logger = structlog.get_logger(__name__)


@HandlerRegistry.register
class PptxHandler(BaseHandler):
    def can_handle(self, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    def extract(self, path: Path) -> Result[RawContent]:
        # Lazy import — see handlers/pdf_handler.py for rationale.
        from core.config import CONFIG

        max_bytes = CONFIG.main.handlers.max_file_size_bytes

        try:
            size = path.stat().st_size
        except OSError as exc:
            logger.warning("pptx.stat.failed", path=str(path), error=str(exc))
            return Failure(
                error=f"PPTX stat failed: {exc}",
                recoverable=False,
                context={"path": str(path)},
            )

        if size > max_bytes:
            logger.warning(
                "pptx.too_large", path=str(path), size=size, limit=max_bytes
            )
            return Failure(
                error=f"PPTX too large: {size} > {max_bytes} bytes",
                recoverable=False,
                context={"path": str(path), "size": size, "limit": max_bytes},
            )

        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                title = ""
                texts: list[str] = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    shape_text = shape.text_frame.text.strip()
                    if not shape_text:
                        continue
                    if (
                        shape.is_placeholder
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    ):
                        title = shape_text
                    else:
                        texts.append(shape_text)
                header = f'[Slide {i}: "{title}"]' if title else f"[Slide {i}]"
                if texts:
                    slides.append(header + "\n" + "\n".join(texts))
                elif title:
                    slides.append(header)
            text = "\n\n".join(slides).strip()
        except Exception as exc:
            return Failure(
                error=f"PPTX read failed: {exc}",
                recoverable=False,
                context={"path": str(path)},
            )
        return Success(RawContent(text=text, source_path=path, is_md=False))
